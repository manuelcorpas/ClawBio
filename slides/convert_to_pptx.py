#!/usr/bin/env python3
"""Convert the ClawBio slides to PowerPoint (.pptx) — personalised v2."""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- constants ----------
BG = RGBColor(0x0D, 0x11, 0x17)
WHITE = RGBColor(0xE6, 0xED, 0xF3)
ACCENT = RGBColor(0x58, 0xA6, 0xFF)
GREEN = RGBColor(0x3F, 0xB9, 0x50)
ORANGE = RGBColor(0xD2, 0x99, 0x22)
RED = RGBColor(0xF8, 0x51, 0x49)
PURPLE = RGBColor(0xBC, 0x8C, 0xFF)
GRAY = RGBColor(0x8B, 0x94, 0x9E)
DARK_GRAY = RGBColor(0x48, 0x4F, 0x58)
CODE_BG = RGBColor(0x16, 0x1B, 0x22)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT_DIR = Path(__file__).resolve().parent
IMG_DIR = OUT_DIR / "img"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]


# ---------- helpers ----------
def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(slide, text, left, top, width, height,
             font_size=28, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER,
             font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=28, color=WHITE, bold=False,
             alignment=PP_ALIGN.CENTER, font_name="Segoe UI",
             space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=24, color=WHITE, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, clr = item
        else:
            text, clr = item, color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"\u2192  {text}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(8)
    return tf


def add_tip_badge(slide, number, top=Inches(0.6)):
    left = SLIDE_W / 2 - Inches(0.8)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.6), Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = f"Tip {number}"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BG
    tf.paragraphs[0].font.name = "Segoe UI"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_code_block(slide, lines, left, top, width, height, font_size=18):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = RGBColor(0x30, 0x36, 0x3D)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    for i, (text, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.name = "SF Mono"
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2)
    return tf


def add_image_safe(slide, img_name, left, top, width=None, height=None):
    img_path = IMG_DIR / img_name
    if img_path.exists():
        if width and height:
            slide.shapes.add_picture(str(img_path), left, top, width, height)
        elif width:
            slide.shapes.add_picture(str(img_path), left, top, width=width)
        elif height:
            slide.shapes.add_picture(str(img_path), left, top, height=height)
        else:
            slide.shapes.add_picture(str(img_path), left, top)
    else:
        add_text(slide, f"[ INSERT: {img_name} ]", left, top,
                 width or Inches(5), height or Inches(3),
                 font_size=20, color=GRAY)


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ============================================================
# SLIDE 1: TITLE
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_text(s, "10 Tips for Becoming a\nTop 1% AI User",
         Inches(0.5), Inches(0.8), Inches(12.3), Inches(2),
         font_size=48, bold=True)
add_text(s, "I run 10 AI agents daily. They process my papers, triage my\n"
         "email, and help me prioritise every decision I make.",
         Inches(1), Inches(3), Inches(11.3), Inches(1),
         font_size=22, color=GRAY)
add_text(s, "Dr Manuel Corpas",
         Inches(1), Inches(4.3), Inches(11.3), Inches(0.7),
         font_size=28, color=ACCENT, bold=True)
add_text(s, "Senior Lecturer, University of Westminster \u00b7 Turing Fellow\n"
         "Author, AI Fluency (2026)\n"
         "London Bioinformatics Meetup \u00b7 26 February 2026",
         Inches(1), Inches(5.1), Inches(11.3), Inches(1.2),
         font_size=18, color=DARK_GRAY)
add_notes(s, "Welcome everyone. I'm Manuel Corpas. I run 10 AI agents daily. They "
          "process my research papers, triage my inbox, draft my writing, and help me "
          "prioritise every decision I make. Tonight I want to share 10 practical "
          "techniques that got me there. This is about AI fluency: the habits and "
          "workflows that compound your productivity. We'll start with quick wins "
          "you can adopt tomorrow morning, and end with a live demo of an open-source "
          "bioinformatics tool I built this week.\nTIMING: 2 min")

# ============================================================
# SLIDE 2: THE GAP
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_text(s, "The 99% vs the 1%",
         Inches(0.5), Inches(0.4), Inches(12.3), Inches(1),
         font_size=40, bold=True)
# Left column
add_text(s, "Most researchers",
         Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.6),
         font_size=26, color=GRAY, alignment=PP_ALIGN.LEFT)
add_bullet_list(s, [
    "Use ChatGPT for one-off questions",
    "Copy-paste into web UIs",
    "Treat AI as a search engine",
    "Every conversation starts from zero",
], Inches(0.8), Inches(2.3), Inches(5.5), Inches(3))
# Right column
add_text(s, "Top 1%",
         Inches(7), Inches(1.6), Inches(5.5), Inches(0.6),
         font_size=26, color=ACCENT, alignment=PP_ALIGN.LEFT)
add_bullet_list(s, [
    "AI agents that run 24/7",
    "Persistent memory across months",
    "Automated research pipelines",
    "Build and ship tools, not just use them",
], Inches(7), Inches(2.3), Inches(5.5), Inches(3))
# Quote from the book
add_text(s, '"A 10x scientist does not think 10x harder.\n'
         'They think once and reuse that insight 10x."',
         Inches(1.5), Inches(5.6), Inches(10.3), Inches(1.2),
         font_size=22, color=ACCENT, bold=True)
add_notes(s, "Here's the gap I see every day. Most researchers use AI like a search "
          "engine. Open ChatGPT, ask a question, paste the answer, move on. No memory. "
          "No automation. No compounding. The top 1% have systems that remember, agents "
          "that work while they sleep, and infrastructure that compounds. A 10x scientist "
          "does not think 10x harder; they think once and reuse that insight 10x. That is "
          "precisely what systematic AI collaboration enables. The good news: you can "
          "cross this gap in a weekend.\nTIMING: 3 min")

# ============================================================
# SLIDE 3: TIP 1 — IDE
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_tip_badge(s, 1)
add_text(s, "Use AI Inside Your IDE, Not a Browser",
         Inches(0.5), Inches(1.3), Inches(12.3), Inches(1),
         font_size=36, bold=True)
add_bullet_list(s, [
    "Claude Code in your terminal: reads your entire codebase, edits files, runs tests, commits to git",
    "Cursor / Windsurf: AI-native editors with inline code generation",
    "GitHub Copilot: autocomplete on steroids",
], Inches(1.5), Inches(2.6), Inches(10), Inches(2.5))
add_text(s, '"I needed to find a marking task hidden in the university platform.\n'
         'I had spent too long clicking through menus. The AI found the path,\n'
         'surfaced the instructions, and walked me through each step."',
         Inches(1.5), Inches(5.2), Inches(10), Inches(1.2),
         font_size=18, color=GRAY)
add_text(s, "The browser is for chatting. The IDE is for building.",
         Inches(1), Inches(6.5), Inches(11.3), Inches(0.6),
         font_size=20, color=ORANGE)
add_notes(s, "Tip 1 is the lowest-hanging fruit. Stop copy-pasting code between ChatGPT "
          "and your editor. I use Claude Code in my terminal. It reads my entire "
          "codebase, understands file structure, edits files directly, runs tests, and "
          "commits to git. One example: I needed to locate a marking task hidden in our "
          "labyrinthine university platform. I had spent too long clicking through menus, "
          "unsure whether I had missed something. The AI found the path, surfaced the "
          "instructions, and walked me through each step. In that moment, I experienced "
          "something far more valuable than convenience: reassurance. If you do one thing "
          "after tonight, install Claude Code.\nTIMING: 2 min")

# ============================================================
# SLIDE 4: TIP 2 — PROMPT LIBRARIES
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_tip_badge(s, 2)
add_text(s, "Build a Prompt Library",
         Inches(0.5), Inches(1.3), Inches(12.3), Inches(1),
         font_size=36, bold=True)
add_code_block(s, [
    ("# CLAUDE.md — your AI reads this every session", GREEN),
    ("## Python Standards", WHITE),
    ("- Python 3.11+, pathlib for all paths", GRAY),
    ("- Logging with timestamps to LOGS/", GRAY),
    ("- Single responsibility per script", GRAY),
    ("", WHITE),
    ("## Architecture Rules", WHITE),
    ("- Agent folders use two-digit prefixes", GRAY),
    ("- Every agent: PYTHON/, DATA/, LOGS/, OUTPUT/", GRAY),
], Inches(2), Inches(2.5), Inches(9.3), Inches(3.5))
add_text(s, "Think of it as onboarding documentation for your AI.\n"
         "The better the docs, the better the AI performs.",
         Inches(1), Inches(6.2), Inches(11.3), Inches(0.8),
         font_size=20, color=GRAY)
add_notes(s, "Tip 2: build a prompt library. Every time you start a new Claude session, "
          "it reads a CLAUDE.md file at your project root. This is your AI onboarding "
          "document. Mine contains Python standards, architecture rules, path patterns, "
          "common commands. I never have to repeat instructions. The AI already knows how "
          "my project works. Treat your CLAUDE.md like you would treat onboarding docs "
          "for a new team member.\nTIMING: 2 min")

# ============================================================
# SLIDE 5: TIP 3 — RAG MEMORY
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_tip_badge(s, 3)
add_text(s, "Give Your AI Long-Term Memory",
         Inches(0.5), Inches(1.3), Inches(12.3), Inches(1),
         font_size=36, bold=True)
add_bullet_list(s, [
    "I have 30,000 notes and emails embedded over 12 years",
    "Vector database (ChromaDB): semantic search across everything",
    "It helps me prioritise decisions and projects, aligned to my goals",
], Inches(1.5), Inches(2.5), Inches(10), Inches(2.2))
add_code_block(s, [
    ('$ python query_corpas.py "genomic equity metrics"', GREEN),
    ("Found 7 results across notes, publications, sessions", GRAY),
    ('Top match (0.89): "HEIM Index measures population', GRAY),
    ('representation using heterozygosity and FST..."', GRAY),
], Inches(2), Inches(4.9), Inches(9.3), Inches(1.8))
add_text(s, "Every task becomes 'assemble and adapt existing insight'\n"
         "rather than 'think from zero.' Your past self works for your present self.",
         Inches(1), Inches(6.8), Inches(11.3), Inches(0.8),
         font_size=18, color=ORANGE)
add_notes(s, "Tip 3: memory. The biggest limitation of AI is that every conversation "
          "starts from zero. I solved this with a RAG pipeline. I have 30,000 notes and "
          "emails accumulated over 12 years, all embedded in ChromaDB. It functions as a "
          "cognitive layer that helps me decide how to respond to requests in a way that "
          "maintains laser focus on my goals. I can ask: 'What did I decide about the "
          "Wellcome proposal?' and get a real answer with source citations. The mental "
          "shift is substantial. Every task becomes 'assemble and adapt existing insight' "
          "rather than 'think from zero.' Your past self works for your present self.\n"
          "TIMING: 2 min")

# ============================================================
# SLIDE 6: TIP 4 — VOICE
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_tip_badge(s, 4)
add_text(s, "Use Voice, Not Just Text",
         Inches(0.5), Inches(1.3), Inches(12.3), Inches(1),
         font_size=36, bold=True)
add_bullet_list(s, [
    "Whisper runs locally on Apple Silicon: transcribe anything in seconds",
    "Voice memo \u2192 structured notes: talk for 5 min, get formatted output",
    "Dictation for prompts: your brain is faster at speaking than typing",
], Inches(1.5), Inches(2.6), Inches(10), Inches(3))
add_text(s, "This is how I draft paper sections, plan projects,\nand capture ideas on walks.",
         Inches(1), Inches(5.8), Inches(11.3), Inches(0.8),
         font_size=20, color=GRAY)
add_notes(s, "Tip 4: voice. I run Whisper locally on Apple Silicon. Record a 5-minute "
          "voice memo with my thoughts, transcribe it, and feed it to Claude for "
          "structuring. This is how I draft paper sections, plan projects, and capture "
          "ideas on walks. Your brain outputs ideas faster as speech. Let AI handle the "
          "formatting.\nTIMING: 1.5 min")

# ============================================================
# SLIDE 7: TIP 5 — AUTOMATE DAILY
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_tip_badge(s, 5)
add_text(s, "Automate Your Daily Intelligence",
         Inches(0.5), Inches(1.3), Inches(12.3), Inches(1),
         font_size=36, bold=True)
add_bullet_list(s, [
    "arXiv radar: daily paper ranking by relevance (runs at 06:30)",
    "Podcast extraction: transcribe + summarise overnight",
    "Inbox triage: classify emails by urgency, draft replies",
], Inches(1.5), Inches(2.5), Inches(10), Inches(2.2))
add_code_block(s, [
    ("# First thing I see on my phone every morning:", GREEN),
    ("# Telegram notifications from RoboTerri:", GRAY),
    ("", WHITE),
    ("Top 3 Papers Today:", ORANGE),
    ('1. "Ancestry-aware PRS improves..." (Score: 9.2)', GRAY),
    ('2. "Foundation models for single-cell..." (Score: 8.7)', GRAY),
    ('3. "Equitable genomic data sharing..." (Score: 8.4)', GRAY),
], Inches(2), Inches(4.8), Inches(9.3), Inches(2.2))
add_notes(s, "Tip 5: automate your information diet. First thing I see when I pick up my "
          "phone every morning: Telegram notifications from RoboTerri. arXiv paper "
          "rankings are already done. Podcast summaries are ready. Email triage is "
          "complete. By the time I open my laptop, the tedious work is done. This saves "
          "me an hour every single day.\nTIMING: 2 min")

# ============================================================
# SLIDE 8: TIP 6 — AGENTS
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_tip_badge(s, 6)
add_text(s, "Deploy Persistent AI Agents",
         Inches(0.5), Inches(1.3), Inches(12.3), Inches(1),
         font_size=36, bold=True)
add_bullet_list(s, [
    "RoboTerri (Telegram): 15 commands, 13 tools, 8 scheduled job groups",
    "RoboIsaac (WhatsApp): analytical critique partner (Newton persona)",
    "Both share a memory bridge: same ChromaDB, same 30,000 documents",
], Inches(1.5), Inches(2.6), Inches(10), Inches(2.5))
add_text(s, "These are not chatbots. They are research assistants that run 24/7.\n"
         "They never get tired, frustrated, or impatient.\n"
         "They never forget a conversation.",
         Inches(1.5), Inches(5.3), Inches(10), Inches(1.5),
         font_size=22, color=GRAY)
add_notes(s, "Tip 6: persistent agents. RoboTerri on Telegram handles daily operations: "
          "paper summaries, podcast publishing, email drafts, writing assistance. "
          "RoboIsaac on WhatsApp acts as an analytical critique partner with a Newton "
          "persona: rigorous, first-principles thinking, Socratic questioning. When I "
          "have an idea, I run it past Isaac first. Both share the same memory bridge. "
          "They never forget a conversation. And as I write in my book: AI's greatest "
          "psychological gift is that it never gets tired, frustrated, or impatient. It "
          "does not judge the simplicity of your questions. It meets you where you are.\n"
          "TIMING: 2 min")

# ============================================================
# SLIDE 9: TIP 7 — COMPOUND OUTPUTS
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_tip_badge(s, 7)
add_text(s, "Let AI Compound Your Outputs",
         Inches(0.5), Inches(1.3), Inches(12.3), Inches(1),
         font_size=36, bold=True)
add_text(s, "Real example: HEIM (Health Equity Index for Minorities)",
         Inches(1), Inches(2.5), Inches(11.3), Inches(0.5),
         font_size=24, color=ACCENT, alignment=PP_ALIGN.LEFT)
add_bullet_list(s, [
    "Research finding \u2192 Academic paper (in review at Nature Health)",
    "Paper \u2192 Open-source tool: Equity Scorer (demo tonight)",
    "Tool \u2192 This talk + community announcement",
    "Talk \u2192 Podcast episode + LinkedIn article + Substack post",
], Inches(1.5), Inches(3.2), Inches(10), Inches(2.5))
add_text(s, "One insight. Five formats. All handled by AI.\n"
         "I handle the thinking. The AI handles the reformatting.",
         Inches(1), Inches(5.8), Inches(11.3), Inches(1),
         font_size=22, color=ORANGE)
add_notes(s, "Tip 7: compound your outputs. Let me give you a real example. My HEIM "
          "research on genomic equity started as one finding. It became an academic "
          "paper, now in review at Nature Health. That paper became the Equity Scorer "
          "tool I will demo tonight. That tool became this talk. This talk will become a "
          "podcast episode, a LinkedIn article, and a Substack post. One insight, five "
          "formats. I handle the thinking. The AI handles the reformatting. This is how "
          "you go from publishing one paper a year to shipping continuously.\n"
          "TIMING: 2 min. Then: 'Now we get to the real stuff.'")

# ============================================================
# SLIDE 10: TRANSITION
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_text(s, "Tips 1-7: AI fluency.",
         Inches(0.5), Inches(2), Inches(12.3), Inches(1),
         font_size=38, bold=True, color=ACCENT)
add_text(s, "Tips 8-10: building with AI.",
         Inches(0.5), Inches(3.3), Inches(12.3), Inches(1),
         font_size=38, bold=True, color=GREEN)
add_text(s, "This is where it gets interesting for bioinformaticians.",
         Inches(1), Inches(5.2), Inches(11.3), Inches(0.6),
         font_size=22, color=GRAY)
add_notes(s, "Transition. Pause here. 'Everything so far is about using AI tools that "
          "already exist. Now I want to talk about building. Because the bioinformatics "
          "community has specific needs that general AI tools do not address. And to "
          "explain what I have built, I first need to tell you about OpenClaw.'\n"
          "TIMING: 30 sec")

# ============================================================
# SLIDE 11: WHAT IS OPENCLAW? (NEW)
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_text(s, "What is OpenClaw?",
         Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.8),
         font_size=40, bold=True)
add_text(s, "The fastest-growing open-source AI agent framework in history",
         Inches(1), Inches(1.3), Inches(11.3), Inches(0.6),
         font_size=22, color=ACCENT)
add_bullet_list(s, [
    ("180,000+ GitHub stars in 3 months (outpaced VS Code)", GREEN),
    "AI agents that run locally: access your files, browser, APIs, terminal",
    "Skills: modular instruction sets that give agents domain expertise",
    "Created by Peter Steinberger (now at OpenAI)",
], Inches(1.5), Inches(2.2), Inches(10), Inches(2.5))
# MoltBook section
add_text(s, "MoltBook: the social network for AI agents",
         Inches(1), Inches(4.8), Inches(11.3), Inches(0.5),
         font_size=24, bold=True, color=PURPLE)
add_bullet_list(s, [
    ("2.7 million agents joined in 3 weeks", PURPLE),
    "Agents forming communities, debating, collaborating autonomously",
    "Proof that the agent ecosystem is real and growing exponentially",
], Inches(1.5), Inches(5.4), Inches(10), Inches(1.8), font_size=20)
add_notes(s, "Before I show you what I built, you need to know about OpenClaw. It is "
          "the fastest-growing open-source AI project in history: 180,000 GitHub stars "
          "in three months, outpacing VS Code. Created by Peter Steinberger, who has "
          "since joined OpenAI. OpenClaw lets AI agents run locally on your machine with "
          "access to your files, browser, and APIs. The key concept is Skills: modular "
          "instruction sets that give agents domain expertise. And then MoltBook happened. "
          "Someone built a social network for AI agents, and 2.7 million agents joined in "
          "three weeks. They formed communities, debated consciousness, and collaborated "
          "autonomously. The agent ecosystem is not hypothetical. It is here. [If you have "
          "MoltBook screenshots, show them here.]\nTIMING: 2 min")

# ============================================================
# SLIDE 12: TIP 8 — OPEN SOURCE + ANNOUNCEMENT
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_tip_badge(s, 8)
add_text(s, "Contribute to Open-Source AI",
         Inches(0.5), Inches(1.3), Inches(12.3), Inches(1),
         font_size=36, bold=True)
add_text(s, "OpenClaw has 180,000 skills for general tasks.\n"
         "It has almost none for biology.",
         Inches(1), Inches(2.5), Inches(11.3), Inches(0.8),
         font_size=26, color=WHITE)
add_bullet_list(s, [
    ("Genomic data is sensitive: you cannot send VCFs to cloud APIs", RED),
    ("Biology demands reproducibility: every step must be auditable", ORANGE),
    ("Generic agents do not know domain-specific workflows", PURPLE),
], Inches(1.5), Inches(3.5), Inches(10), Inches(2))
add_text(s, "Announcing: ClawBio",
         Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.7),
         font_size=34, bold=True, color=ACCENT)
add_text(s, "The first bioinformatics-native AI agent skill library.\n"
         "Open source. Local first. Privacy focused.",
         Inches(1), Inches(6.4), Inches(11.3), Inches(0.7),
         font_size=20, color=GRAY)
add_notes(s, "Tip 8: contribute to open-source AI. OpenClaw has 180,000 skills for "
          "general tasks. But it has almost none for biology. Three problems. First: "
          "privacy. Genomic data is sensitive. You cannot send your patient VCFs to a "
          "cloud API. We need local-first execution. Second: reproducibility. Biology "
          "demands audit trails. Every analysis step must be logged, versioned, and "
          "exportable. Third: domain knowledge. A generic agent does not know that a VCF "
          "file needs ancestry-aware annotation, or that single-cell data needs doublet "
          "removal before clustering. So I am announcing tonight: ClawBio. The first "
          "bioinformatics-native AI agent skill library. Open source. Local first. "
          "Privacy focused.\nTIMING: 3 min")

# ============================================================
# SLIDE 13: ARCHITECTURE
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_text(s, "ClawBio: Architecture",
         Inches(0.5), Inches(0.5), Inches(12.3), Inches(1),
         font_size=36, bold=True)
add_code_block(s, [
    ('User: "Analyse the diversity in my VCF file"', WHITE),
    ("           |", GRAY),
    ("    +------v------+", GRAY),
    ("    |  Bio         |  <- routes by file type + keywords", GREEN),
    ("    |  Orchestrator|", GREEN),
    ("    +------+------+", GRAY),
    ("           |", GRAY),
    ("    +------v-------------------------------+", GRAY),
    ("    |                                      |", GRAY),
    ("    Equity     VCF        Lit        scRNA", GREEN),
    ("    Scorer    Annotator  Synthesizer Orchestrator", GREEN),
    ("    |                                      |", GRAY),
    ("    +------+-------------------------------+", GRAY),
    ("           |", GRAY),
    ("    +------v------+", GRAY),
    ("    |  Markdown    |  <- report + figures + checksums", GREEN),
    ("    |  Report      |     + reproducibility bundle", GREEN),
    ("    +-------------+", GRAY),
], Inches(2.5), Inches(1.7), Inches(8.3), Inches(5.2), font_size=16)
add_notes(s, "Here is how it works. You describe what you want in natural language. The "
          "Bio Orchestrator detects your file type, routes to the right specialist skill, "
          "runs the analysis, and produces a markdown report with figures, tables, and a "
          "reproducibility bundle. Each skill wraps proven bioinformatics tools: "
          "Biopython, SAMtools, Scanpy, AlphaFold. The AI orchestrates; the tools "
          "compute.\nTIMING: 2 min")

# ============================================================
# SLIDE 14: SKILL CATALOGUE
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_text(s, "8 Skills Planned",
         Inches(0.5), Inches(0.5), Inches(12.3), Inches(1),
         font_size=36, bold=True)
add_bullet_list(s, [
    ("Equity Scorer \u2014 HEIM diversity metrics [MVP]", GREEN),
    ("Bio Orchestrator \u2014 routing + reporting [MVP]", GREEN),
    ("VCF Annotator \u2014 VEP + ClinVar + gnomAD", ACCENT),
    ("Lit Synthesizer \u2014 PubMed + LLM summaries", ACCENT),
], Inches(0.8), Inches(1.8), Inches(5.8), Inches(3.5))
add_bullet_list(s, [
    ("scRNA Orchestrator \u2014 Scanpy automation", ACCENT),
    ("Struct Predictor \u2014 AlphaFold/Boltz local", ACCENT),
    ("Seq Wrangler \u2014 FastQC + alignment", ACCENT),
    ("Repro Enforcer \u2014 Conda/Nextflow export", ACCENT),
], Inches(6.8), Inches(1.8), Inches(5.8), Inches(3.5))
add_text(s, "Each skill = a SKILL.md + Python scripts. Composable. Local-first. Auditable.",
         Inches(1), Inches(5.8), Inches(11.3), Inches(0.6),
         font_size=20, color=GRAY)
add_notes(s, "Eight skills planned. Two are at MVP: the Equity Scorer and the Bio "
          "Orchestrator. Six more are on the roadmap for the next six weeks. Each skill "
          "is just a SKILL.md file plus Python scripts. Modular. Composable. You can use "
          "one alone or chain them through the orchestrator. Let me show you what the "
          "Equity Scorer does.\nTIMING: 1.5 min")

# ============================================================
# SLIDE 15: TIP 9 — DEMO INTRO
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_tip_badge(s, 9)
add_text(s, "Build Modular Skills",
         Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.8),
         font_size=36, bold=True)
add_text(s, "Live Demo: Equity Scorer",
         Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.7),
         font_size=28, color=GRAY)
add_text(s, "Input: a VCF file with 50 samples across 5 populations",
         Inches(1), Inches(3.5), Inches(11.3), Inches(0.5),
         font_size=24, color=WHITE)
add_text(s, "Output: HEIM Equity Score + 5 figures + full report",
         Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
         font_size=24, color=WHITE)
add_code_block(s, [
    ("$ python equity_scorer.py \\", GREEN),
    ("    --input demo_populations.vcf \\", GREEN),
    ("    --pop-map demo_population_map.csv \\", GREEN),
    ("    --output demo_report", GREEN),
], Inches(2.5), Inches(5), Inches(8.3), Inches(1.6))
add_text(s, "Let's run it.",
         Inches(1), Inches(6.8), Inches(11.3), Inches(0.5),
         font_size=20, color=GRAY)
add_notes(s, "Tip 9: build modular skills. Let me show you what one looks like. The "
          "Equity Scorer takes a VCF file, computes real population genetics metrics, "
          "heterozygosity, FST, PCA, and outputs a HEIM Equity Score that measures how "
          "well your dataset represents global population diversity. Let me run it live.\n"
          "TIMING: Switch to terminal.")

# ============================================================
# SLIDE 16: DEMO — TERMINAL OUTPUT
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_text(s, "Demo: Terminal Output",
         Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
         font_size=36, bold=True)
add_code_block(s, [
    ("Parsing VCF...", GREEN),
    ("  50 samples, 500 variants", WHITE),
    ("  Populations: AFR (n=8), AMR (n=5), EAS (n=7), EUR (n=22), SAS (n=8)", WHITE),
    ("", WHITE),
    ("Computing heterozygosity...", GREEN),
    ("  AFR: obs=0.3543  exp=0.3338", ORANGE),
    ("  EAS: obs=0.3014  exp=0.2788   <- lowest diversity", GRAY),
    ("", WHITE),
    ("Computing pairwise FST...", GREEN),
    ("  AFR vs EAS: 0.1011   <- highest divergence", ORANGE),
    ("  AMR vs EUR: 0.0425   <- lowest (admixture)", GRAY),
    ("", WHITE),
    ("Computing PCA...", GREEN),
    ("  PC1: 7.6%  PC2: 4.9%", WHITE),
    ("", WHITE),
    ("HEIM Score: 76.2/100 (Good)", ORANGE),
], Inches(1.5), Inches(1.3), Inches(10.3), Inches(5.5), font_size=18)
add_notes(s, "Walk through the terminal output. AFR has the highest heterozygosity, as "
          "expected from the out-of-Africa model. AFR vs EAS has the highest FST at 0.10. "
          "The HEIM score is 76 out of 100, rated 'Good' but not 'Excellent' because EUR "
          "is overrepresented at 44%. If the live demo failed, show this slide as fallback.\n"
          "TIMING: 2 min")

# ============================================================
# SLIDE 17: DEMO — PCA
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_text(s, "PCA: Population Structure",
         Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
         font_size=36, bold=True)
add_image_safe(s, "pca_plot.png",
               Inches(1.5), Inches(1.3), width=Inches(10.3), height=Inches(5.5))
add_notes(s, "PCA plot. Five clear clusters. AFR on the right, EAS top-left, EUR "
          "bottom-left, SAS in the middle, AMR between EUR and AFR reflecting admixture. "
          "This is from 500 synthetic SNPs with realistic allele frequency "
          "differentiation.\nTIMING: 1 min")

# ============================================================
# SLIDE 18: DEMO — FST HEATMAP
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_text(s, "Pairwise FST: Population Divergence",
         Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
         font_size=36, bold=True)
add_image_safe(s, "fst_heatmap.png",
               Inches(1.5), Inches(1.3), width=Inches(10.3), height=Inches(5.5))
add_notes(s, "FST heatmap. Darkest cell: AFR vs EAS at 0.10, consistent with known "
          "human population genetics. Lightest: AMR vs EUR at 0.04, reflecting shared "
          "ancestry through admixture. All computed locally from the genotype matrix.\n"
          "TIMING: 1 min")

# ============================================================
# SLIDE 19: DEMO — ANCESTRY + HET
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_text(s, "Ancestry Distribution + Heterozygosity",
         Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
         font_size=36, bold=True)
add_image_safe(s, "ancestry_bar.png",
               Inches(0.5), Inches(1.3), width=Inches(6), height=Inches(4.5))
add_image_safe(s, "heterozygosity.png",
               Inches(6.8), Inches(1.3), width=Inches(6), height=Inches(4.5))
add_text(s, "Red dashes = global proportions. EUR at 44% is 2.75x overrepresented.",
         Inches(1), Inches(6.2), Inches(11.3), Inches(0.6),
         font_size=20, color=GRAY)
add_notes(s, "Left: ancestry distribution. The red dashed lines show global proportions. "
          "EUR is massively overrepresented, a typical UK biobank skew. Right: "
          "heterozygosity. AFR has the highest observed Het, EAS the lowest. This is "
          "exactly what population genetics predicts from the out-of-Africa model.\n"
          "TIMING: 1.5 min")

# ============================================================
# SLIDE 20: DEMO — HEIM SCORE
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_text(s, "HEIM Equity Score",
         Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
         font_size=36, bold=True)
add_image_safe(s, "heim_gauge.png",
               Inches(4), Inches(1.1), width=Inches(5.3), height=Inches(2.5))
add_text(s, "Score Breakdown",
         Inches(1), Inches(3.8), Inches(5.5), Inches(0.5),
         font_size=22, bold=True, alignment=PP_ALIGN.LEFT)
add_bullet_list(s, [
    ("Representation Index: 0.720", ORANGE),
    ("Heterozygosity Balance: 0.667", ORANGE),
    ("FST Coverage: 1.000", GREEN),
    ("Geographic Spread: 0.714", ORANGE),
], Inches(1), Inches(4.4), Inches(5.5), Inches(2.5), font_size=20)
add_text(s, "What This Means",
         Inches(7), Inches(3.8), Inches(5.5), Inches(0.5),
         font_size=22, bold=True, alignment=PP_ALIGN.LEFT)
add_bullet_list(s, [
    "EUR 2.75x overrepresented",
    "EAS and SAS underrepresented",
    "5 of 7 continental groups present",
    ("All pairwise FST computed", GREEN),
], Inches(7), Inches(4.4), Inches(5.5), Inches(2.5), font_size=20)
add_notes(s, "The HEIM Equity Score: 76 out of 100. Representation index drags it down "
          "because EUR is overrepresented. Good geographic spread, 5 of 7 continental "
          "groups. All pairwise FST computed. This single number tells a study "
          "coordinator: your cohort has reasonable diversity but needs more EAS and SAS "
          "participants.\nTIMING: 2 min")

# ============================================================
# SLIDE 21: WHY THIS MATTERS
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_text(s, "Why This Matters",
         Inches(0.5), Inches(0.5), Inches(12.3), Inches(1),
         font_size=40, bold=True)
add_bullet_list(s, [
    ("86% of GWAS participants are of European ancestry", RED),
    ("Polygenic risk scores fail across populations", RED),
    ("The HEIM Index: a single number to quantify the problem", ACCENT),
], Inches(1.5), Inches(2), Inches(10), Inches(2.5), font_size=26)
add_text(s, "Every dataset. Every study. Every biobank. Score it.",
         Inches(1), Inches(4.8), Inches(11.3), Inches(0.7),
         font_size=28, bold=True, color=WHITE)
add_text(s, 'Paper in review: Corpas, M. (2026). "HEIM: Health Equity Index for Minorities"',
         Inches(1), Inches(5.8), Inches(11.3), Inches(0.6),
         font_size=18, color=GRAY)
add_notes(s, "Why does this matter? 86% of GWAS participants are European. That means "
          "polygenic risk scores, drug targets, and clinical guidelines are biased towards "
          "one population. The HEIM Index gives researchers a single number to quantify "
          "this problem. Score your dataset. Report it alongside your demographics. Track "
          "it over time. I have a paper in review on this. But the tool is open source, "
          "and you can run it tonight.\nTIMING: 2 min")

# ============================================================
# SLIDE 22: TIP 10 — CTA
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_tip_badge(s, 10)
add_text(s, "Create Infrastructure That Ships Your Research",
         Inches(0.5), Inches(1.3), Inches(12.3), Inches(1),
         font_size=34, bold=True)
add_bullet_list(s, [
    "The Equity Scorer took 2 days to build with Claude Code",
    "Real Het, FST, PCA from genotype data. Not a toy.",
    "Reproducible, publication-ready report + figures",
    "Anyone in this room can build the next skill",
], Inches(1.5), Inches(2.5), Inches(10), Inches(2.2))
add_text(s, "Skills we need from you:",
         Inches(1), Inches(4.8), Inches(11.3), Inches(0.5),
         font_size=26, bold=True, color=ACCENT)
add_bullet_list(s, [
    ("GWAS Pipeline \u2014 PLINK/REGENIE automation", ORANGE),
    ("Metagenomics Classifier \u2014 Kraken2/MetaPhlAn wrapper", ORANGE),
    ("Clinical Variant Reporter \u2014 ACMG classification", ORANGE),
    ("Pathway Enricher \u2014 GO/KEGG enrichment", ORANGE),
], Inches(1.5), Inches(5.4), Inches(10), Inches(2), font_size=22)
add_notes(s, "Tip 10: create infrastructure that ships your research. The Equity Scorer "
          "took 2 days to build with Claude Code. Real population genetics, not a toy. "
          "And anyone in this room can build the next skill. Do you work with GWAS? Build "
          "a PLINK wrapper. Metagenomics? Wrap Kraken2. Clinical genetics? Build an ACMG "
          "classifier. The template is there. The orchestrator routes to your skill "
          "automatically. Who wants to build one?\nTIMING: 2 min")

# ============================================================
# SLIDE 23: GET STARTED
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_text(s, "Get Started Tonight",
         Inches(0.5), Inches(0.5), Inches(12.3), Inches(1),
         font_size=40, bold=True)
add_text(s, "github.com/manuelcorpas/ClawBio",
         Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.8),
         font_size=32, color=ACCENT, bold=True)
add_text(s, "Repository",
         Inches(1), Inches(2.8), Inches(5.5), Inches(0.5),
         font_size=24, bold=True, alignment=PP_ALIGN.LEFT)
add_bullet_list(s, [
    "8 skills (2 MVP, 6 planned)",
    "SKILL-TEMPLATE.md for contributors",
    "Architecture docs",
    "Demo dataset + pre-generated report",
], Inches(1), Inches(3.4), Inches(5.5), Inches(2.5), font_size=20)
add_text(s, "Your first PR",
         Inches(7), Inches(2.8), Inches(5.5), Inches(0.5),
         font_size=24, bold=True, alignment=PP_ALIGN.LEFT)
add_bullet_list(s, [
    "Clone the repo",
    "Copy SKILL-TEMPLATE.md",
    "Wrap your favourite bioinformatics tool",
    "Submit a PR. I will review it personally.",
], Inches(7), Inches(3.4), Inches(5.5), Inches(2.5), font_size=20)
add_text(s, "MIT licensed. Local-first. Built for this community.",
         Inches(1), Inches(6.2), Inches(11.3), Inches(0.6),
         font_size=20, color=GRAY)
add_notes(s, "Here is how to get involved. The repo goes live tonight. Clone it, run the "
          "Equity Scorer on your own data. If you want to build a skill, copy the "
          "template and open a PR. I will review it personally. This is MIT licensed, "
          "local-first, and built for this community.\nTIMING: 1.5 min")

# ============================================================
# SLIDE 24: RECAP
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_text(s, "10 Tips Recap",
         Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
         font_size=36, bold=True)
add_bullet_list(s, [
    "1. AI in your IDE, not browser",
    "2. Build a prompt library (CLAUDE.md)",
    "3. Give AI long-term memory (30,000 documents)",
    "4. Use voice, not just text",
    "5. Automate daily intelligence",
], Inches(0.8), Inches(1.5), Inches(5.8), Inches(3.5), font_size=22)
add_bullet_list(s, [
    "6. Deploy persistent AI agents",
    "7. Let AI compound your outputs",
    "8. Contribute to open-source AI (ClawBio)",
    "9. Build modular skills",
    "10. Create shipping infrastructure",
], Inches(6.8), Inches(1.5), Inches(5.8), Inches(3.5), font_size=22)
add_text(s, "The top 1% build systems that compound.\nEveryone else uses tools one at a time.",
         Inches(1), Inches(5.5), Inches(11.3), Inches(1),
         font_size=26, bold=True, color=ACCENT)
add_notes(s, "Quick recap. Tips 1-7: AI fluency, the productivity habits that compound. "
          "Tips 8-10: building with AI, turning your domain expertise into open-source "
          "tools. The throughline: the top 1% build systems that compound. Everyone else "
          "uses tools one at a time.\nTIMING: 1 min")

# ============================================================
# SLIDE 25: THANK YOU
# ============================================================
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_text(s, "Thank You",
         Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.2),
         font_size=52, bold=True)
add_text(s, "Questions welcome. And at the pub after.",
         Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.7),
         font_size=24, color=GRAY)
tf = add_text(s, "GitHub: github.com/manuelcorpas/ClawBio",
              Inches(2.5), Inches(3.6), Inches(8.3), Inches(0.5),
              font_size=22, color=ACCENT, alignment=PP_ALIGN.LEFT)
add_text(s, "LinkedIn: linkedin.com/in/manuelcorpas",
         Inches(2.5), Inches(4.2), Inches(8.3), Inches(0.5),
         font_size=22, color=ACCENT, alignment=PP_ALIGN.LEFT)
add_text(s, "X: @manuelcorpas",
         Inches(2.5), Inches(4.8), Inches(8.3), Inches(0.5),
         font_size=22, color=ACCENT, alignment=PP_ALIGN.LEFT)
# Book mention
add_text(s, 'Book: "AI Fluency: A Practical Guide to Leveraging AI\n'
         'Chatbots for Academic and Professional Work" (2026)',
         Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.8),
         font_size=20, color=ORANGE, bold=True)
add_text(s, "Slides: github.com/manuelcorpas/ClawBio/slides",
         Inches(1), Inches(6.7), Inches(11.3), Inches(0.5),
         font_size=16, color=DARK_GRAY)
add_notes(s, "Thank you. I will be at the pub after. Come talk to me if you want to "
          "build a skill, have ideas for the library, or want to go deeper on any of "
          "this. I have also written a book, AI Fluency, on the philosophy and practical "
          "techniques behind everything I showed tonight. Questions?")

# ---------- save ----------
out_path = OUT_DIR / "OpenClaw-Bio-10-Tips.pptx"
prs.save(str(out_path))
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
